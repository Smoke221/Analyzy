import magic
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

def detect_file_type(file_path):
    mime = magic.Magic(mime=True)
    file_type = mime.from_file(file_path)
    return file_type

def extract_text(file_path):
    file_type = detect_file_type(file_path)
    if "pdf" in file_type:
        return extract_text_from_pdf(file_path)
    elif "wordprocessingml.document" in file_type:
        return extract_text_from_docx(file_path)
    elif "presentation" in file_type:
        return extract_text_from_ppt(file_path)
    else:
        return "Unsupported file format"

def extract_text_from_pdf(file_path):
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    text = ""
    for para in doc.paragraphs:
        text += para.text
    return text

def extract_text_from_ppt(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

extracted_text = extract_text("samplepptx.pptx")


output_file_path = "extracted_text.txt"
with open(output_file_path, "w", encoding="utf-8") as output_file:
    output_file.write(extracted_text)

print("Extracted text saved to:", output_file_path)
